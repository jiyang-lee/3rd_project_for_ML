"""M1 장애 예측 모델 개발 여정 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
import os

# === Color Palette: Deep Navy + Teal ===
NAVY = RGBColor(0x1A, 0x27, 0x44)
DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x08, 0x91, 0xB2)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
SLATE_MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)


def add_shadow(shape):
    """Add subtle shadow to shape via XML manipulation."""
    from lxml import etree
    spPr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        return
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    effectLst = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
    outerShdw = etree.SubElement(effectLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw',
                                  attrib={'blurRad': '50800', 'dist': '25400', 'dir': '2700000', 'algn': 'tl'})
    srgbClr = etree.SubElement(outerShdw, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                                attrib={'val': '000000'})
    etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', attrib={'val': '15000'})


def set_rounded_rect(shape, radius_emu=100000):
    """Make a rectangle rounded."""
    from lxml import etree
    prstGeom = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            avLst.clear()
        etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd',
                         attrib={'name': 'adj', 'fmla': f'val {radius_emu}'})


def add_card(slide, x, y, w, h, fill=CARD_BG, shadow=True, rounded=True):
    """Add a card (rounded rect with shadow)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        set_rounded_rect(shape)
    if shadow:
        add_shadow(shape)
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=14, color=SLATE, bold=False,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font_name="Calibri", margin=None):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if margin is not None:
        tf.margin_left = Inches(margin)
        tf.margin_right = Inches(margin)
        tf.margin_top = Inches(margin)
        tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_rich_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                  line_spacing=None, space_after=None):
    """Add text box with multiple formatted runs. Each run: (text, size, color, bold, font)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if space_after is not None:
        p.space_after = Pt(space_after)
    for i, run_data in enumerate(runs):
        text, size, color, bold = run_data[:4]
        font_name = run_data[4] if len(run_data) > 4 else "Calibri"
        if i == 0:
            r = p.runs[0] if p.runs else p.add_run()
            r.text = text
        else:
            r = p.add_run()
            r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font_name
    return txBox


def add_bullet_list(slide, x, y, w, h, items, font_size=14, color=SLATE, bold_items=None,
                    space_after=6):
    """Add bulleted list. bold_items: set of indices to bold."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    bold_items = bold_items or set()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.font.bold = i in bold_items
        p.space_after = Pt(space_after)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar',
                                   attrib={'char': '•'})
        buSzPct = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct',
                                    attrib={'val': '100000'})
        indent = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont',
                                   attrib={'typeface': 'Arial'})
    return txBox


def add_table_to_slide(slide, x, y, w, rows_data, col_widths=None, header_color=DEEP_BLUE,
                       font_size=11):
    """Add a formatted table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    row_h = Inches(0.38)
    table_h = row_h * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), table_h)
    table = shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                paragraph.space_after = Pt(0)
                paragraph.space_before = Pt(0)
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = SLATE
                    paragraph.alignment = PP_ALIGN.LEFT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return shape


def dark_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

def light_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

def white_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_slide_number(slide, num, total=22, dark=False):
    color = RGBColor(0x94, 0xA3, 0xB8) if dark else SLATE_MUTED
    add_text_box(slide, 9.0, 5.25, 0.8, 0.3, f"{num} / {total}", font_size=9,
                 color=color, align=PP_ALIGN.RIGHT)


def add_section_tag(slide, text, x=0.5, y=0.15):
    """Add a small section tag at top of slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    set_rounded_rect(shape, 50000)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


# ========================================
# BUILD PRESENTATION
# ========================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_slide_bg(slide)

add_text_box(slide, 0.8, 1.2, 8.4, 1.2,
             "M1 장애 예측 모델", font_size=44, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT)
add_text_box(slide, 0.8, 2.3, 8.4, 0.8,
             "27개 실험의 여정과 최종 구조", font_size=24, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.LEFT)

# Teal accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.3), Inches(2.0), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = TEAL
shape.line.fill.background()

add_text_box(slide, 0.8, 3.6, 8.4, 0.5,
             "HeatGrid  |  지역난방 기계실 예지보전  |  M1 (Manufacturer 1)", font_size=13,
             color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)
add_text_box(slide, 0.8, 4.2, 8.4, 0.4,
             "2026.06", font_size=12, color=SLATE_MUTED, align=PP_ALIGN.LEFT)
add_slide_number(slide, 1, dark=True)


# ========== SLIDE 2: Executive Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.8, 0.3, 8.4, 0.6,
             "결론 먼저", font_size=36, color=WHITE, bold=True)

# Quote box
card = add_card(slide, 0.8, 1.0, 8.4, 0.7, fill=RGBColor(0x0E, 0x34, 0x60), shadow=False, rounded=True)
add_text_box(slide, 1.1, 1.1, 8.0, 0.5,
             '"4분류 모델이 아니라, 계층형 운영 판단 정책이다."',
             font_size=18, color=LIGHT_TEAL, bold=True, align=PP_ALIGN.CENTER)

# Three summary cards
cards_data = [
    ("Fault Gate", "Recall 0.89  |  FPR 0.20", "장애 55건 중 49건 탐지\n정상 35건 중 7건 오탐", GREEN),
    ("Task", "사후 맥락 분류", "예측 대상이 아님\npost-event context", AMBER),
    ("Activity", "보조 신호", "결측 패턴 의존\nmissingness-sensitive", RED),
]
for i, (title, metric, desc, accent) in enumerate(cards_data):
    cx = 0.8 + i * 2.95
    card = add_card(slide, cx, 2.0, 2.7, 2.3, fill=RGBColor(0x0E, 0x34, 0x60), shadow=True, rounded=True)
    # Accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.15), Inches(2.2), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    add_text_box(slide, cx + 0.45, 2.15, 2.0, 0.3, title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, cx + 0.15, 2.55, 2.4, 0.35, metric, font_size=13, color=TEAL, bold=True)
    add_text_box(slide, cx + 0.15, 2.95, 2.4, 1.0, desc, font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

add_slide_number(slide, 2, dark=True)


# ========== SLIDE 3: Problem Statement ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 9.0, 0.6,
             '"고객이 추위를 느끼기 전에 장애를 잡을 수 있는가?"',
             font_size=28, color=SLATE, bold=True)
add_section_tag(slide, "문제 정의", x=0.5, y=0.9)

# Before card
card = add_card(slide, 0.5, 1.4, 4.2, 3.2, fill=RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, 0.7, 1.55, 3.8, 0.35, "현재: 사후 대응", font_size=18, color=RED, bold=True)
add_bullet_list(slide, 0.7, 2.0, 3.8, 2.2, [
    "고객이 추위를 느낀 후 신고",
    "신고 접수 → 현장 출동 → 진단",
    "장애 발생 후 수 시간~수 일 경과",
    "운영자는 장애 인지 불가",
], font_size=12, color=SLATE)

# After card
card = add_card(slide, 5.3, 1.4, 4.2, 3.2, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 5.5, 1.55, 3.8, 0.35, "목표: 선제 대응", font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, 5.5, 2.0, 3.8, 2.2, [
    "기계실 센서 데이터 기반 탐지",
    "장애 3~7일 전 위험 신호 감지",
    "위험도 기반 점검/출동 판단",
    "운영자에게 구조화된 정보 제공",
], font_size=12, color=SLATE)

# Arrow between
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.65), Inches(2.7), Inches(0.7), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = TEAL
arrow.line.fill.background()

add_slide_number(slide, 3)


# ========== SLIDE 4: Data Reality ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 9.0, 0.5, "M1 데이터 현실", font_size=30, color=SLATE, bold=True)
add_text_box(slide, 0.5, 0.8, 9.0, 0.3, "작은 데이터, 불균형 라벨에서 신뢰할 수 있는 모델을 만드는 것이 핵심 과제",
             font_size=13, color=SLATE_MUTED)

# Stat cards
stats = [
    ("35", "정상 이벤트", DEEP_BLUE),
    ("55", "장애 이벤트", RED),
    ("42", "작업(Task)", AMBER),
    ("47", "활동(Activity)", TEAL),
    ("29", "기계실 수", SLATE),
    ("10", "센서 종류", SLATE),
]
for i, (num, label, color) in enumerate(stats):
    col = i % 3
    row = i // 3
    cx = 0.7 + col * 3.1
    cy = 1.4 + row * 1.9
    card = add_card(slide, cx, cy, 2.7, 1.6, fill=WHITE)
    add_text_box(slide, cx + 0.2, cy + 0.2, 2.3, 0.8, num, font_size=48, color=color, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.2, cy + 1.05, 2.3, 0.4, label, font_size=13, color=SLATE_MUTED,
                 align=PP_ALIGN.CENTER)

add_slide_number(slide, 4)


# ========== SLIDE 5: Roadmap ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 9.0, 0.5, "27개 실험의 흐름", font_size=30, color=SLATE, bold=True)
add_text_box(slide, 0.5, 0.8, 9.0, 0.3, "3단계로 나누면", font_size=14, color=SLATE_MUTED)

# Phase boxes
phases = [
    ("Phase 1", "NB 1-13", "Normal vs Fault 분류", "기본 신호 확인 → compact13 발견\n→ Fault Gate 확립",
     DEEP_BLUE, "BA 0.59 → 0.85"),
    ("Phase 2", "NB 14-25", "4분류 시도 + 라벨 발견", "4분류 시도 → 라벨 의미 차이 발견\n→ 계층형 정책 설계",
     TEAL, "라벨 ≠ 같은 종류"),
    ("Phase 3", "NB 27", "최종 Recipe 적용", "13번 pre-event recipe를\nfault gate에 적용",
     GREEN, "최종 확정"),
]
for i, (phase, nb, title, desc, color, tag) in enumerate(phases):
    cx = 0.5 + i * 3.15
    card = add_card(slide, cx, 1.35, 2.9, 3.3, fill=WHITE)
    # Phase label
    plbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + 0.15), Inches(1.5), Inches(1.2), Inches(0.28))
    plbl.fill.solid()
    plbl.fill.fore_color.rgb = color
    plbl.line.fill.background()
    set_rounded_rect(plbl, 50000)
    tf = plbl.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, cx + 0.15, 1.9, 2.6, 0.25, nb, font_size=11, color=SLATE_MUTED)
    add_text_box(slide, cx + 0.15, 2.15, 2.6, 0.35, title, font_size=15, color=SLATE, bold=True)
    add_text_box(slide, cx + 0.15, 2.55, 2.6, 1.2, desc, font_size=12, color=SLATE_MUTED)
    # Tag at bottom
    add_text_box(slide, cx + 0.15, 3.9, 2.6, 0.3, tag, font_size=12, color=color, bold=True,
                 align=PP_ALIGN.CENTER)

# Arrows between phases
for i in range(2):
    ax = 3.4 + i * 3.15
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(2.7), Inches(0.3), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()

add_slide_number(slide, 5)


# ========== SLIDE 6: Baseline ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 1")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "출발점: 50건, 70 Feature (NB 3)", font_size=28, color=SLATE, bold=True)

# Chart: bar chart
chart_data = {
    'categories': ['Dummy Classifier', 'Logistic (balanced)'],
    'BA': [0.50, 0.59],
    'Recall': [0.00, 0.60],
}
from pptx.chart.data import CategoryChartData
cd = CategoryChartData()
cd.categories = chart_data['categories']
cd.add_series('Balanced Accuracy', chart_data['BA'])
cd.add_series('Recall', chart_data['Recall'])

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5), cd
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = SLATE_MUTED
plot = chart.plots[0]
plot.gap_width = 120
series0 = plot.series[0]
series0.format.fill.solid()
series0.format.fill.fore_color.rgb = DEEP_BLUE
series1 = plot.series[1]
series1.format.fill.solid()
series1.format.fill.fore_color.rgb = TEAL

# Right side: interpretation
card = add_card(slide, 6.3, 1.2, 3.2, 3.5, fill=WHITE)
add_text_box(slide, 6.5, 1.4, 2.8, 0.3, "해석", font_size=16, color=SLATE, bold=True)
add_bullet_list(slide, 6.5, 1.8, 2.8, 2.5, [
    "50건의 소규모 데이터",
    "70개 기본 feature (10센서 x 7통계)",
    "약한 신호는 존재",
    "하지만 아직 쓸 수 없는 수준",
    "→ Feature 확장/압축 필요",
], font_size=11, color=SLATE)

# Decision tag
add_text_box(slide, 6.5, 4.0, 2.8, 0.4, "결정: Feature 확장 탐색 시작",
             font_size=11, color=TEAL, bold=True)

add_slide_number(slide, 6)


# ========== SLIDE 7: Feature expansion → compact ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 1")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "13개 Feature로 154개보다 낫다 (NB 8-9)",
             font_size=28, color=SLATE, bold=True)

# Three comparison cards
comp_data = [
    ("base 70", "BA 0.62", "FPR 0.26", "70개 기본 feature\n10센서 x 7통계", SLATE_MUTED),
    ("expanded 154", "BA 0.65 (+2.8%)", "FPR 0.20", "+84개 파생 신호\nsupply_error, return_gap 등", AMBER),
    ("compact 13", "BA 0.83 (+18%)", "FPR 0.06", "핵심 13개만 선택\n92% feature 감소", GREEN),
]
for i, (name, ba, fpr, desc, accent) in enumerate(comp_data):
    cx = 0.5 + i * 3.15
    card = add_card(slide, cx, 1.2, 2.9, 3.4, fill=WHITE)
    add_text_box(slide, cx + 0.2, 1.35, 2.5, 0.3, name, font_size=18, color=accent, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.2, 1.75, 2.5, 0.35, ba, font_size=20, color=SLATE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.2, 2.15, 2.5, 0.25, fpr, font_size=13, color=SLATE_MUTED,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.2, 2.6, 2.5, 1.2, desc, font_size=12, color=SLATE_MUTED,
                 align=PP_ALIGN.CENTER)

# Key insight at bottom
add_text_box(slide, 0.5, 4.8, 9.0, 0.4,
             "핵심: Feature를 늘리는 것보다 핵심을 고르는 것이 낫다",
             font_size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

add_slide_number(slide, 7)


# ========== SLIDE 8: Why compact13 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 1")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "compact13이 핵심인 이유", font_size=28, color=SLATE, bold=True)

# Left: feature list
card = add_card(slide, 0.5, 1.2, 5.0, 3.8, fill=WHITE)
add_text_box(slide, 0.7, 1.35, 4.6, 0.3, "선택된 13개 Feature 카테고리", font_size=15, color=SLATE, bold=True)

features = [
    "온도 관련: supply/return temp 통계 (mean, std, delta)",
    "온도차: return_gap (공급-환수 온도차)",
    "유량: flow rate 변동 (std, last_minus_first)",
    "압력: pressure delta 변화",
    "시간 변화: 12h/6h/1d 구간 delta",
    "결측: missing_rate (데이터 품질 지표)",
]
add_bullet_list(slide, 0.7, 1.75, 4.6, 2.8, features, font_size=12, color=SLATE)

# Right: why it works
card = add_card(slide, 5.8, 1.2, 3.7, 3.8, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 6.0, 1.35, 3.3, 0.3, "왜 더 적은데 더 좋은가?", font_size=15, color=GREEN, bold=True)
add_bullet_list(slide, 6.0, 1.75, 3.3, 2.8, [
    "노이즈 feature 제거 → 과적합 감소",
    "물리적으로 의미 있는 신호에 집중",
    "Fold 간 Jaccard 34~40% 안정성",
    "해석 가능성 대폭 향상",
    "운영자가 \"왜 위험인지\" 설명 가능",
], font_size=12, color=SLATE)

add_slide_number(slide, 8)


# ========== SLIDE 9: Safety Validation ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 1")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "안전성 검증 (NB 10-12)", font_size=28, color=SLATE, bold=True)
add_text_box(slide, 0.5, 0.9, 9.0, 0.3, "compact13은 안전한가?", font_size=14, color=SLATE_MUTED)

# Checklist cards
checks = [
    ("Data Leakage 검증", "Pass", "train/test에 같은 기계실 없음\nSubstation Group CV 적용", GREEN),
    ("Hard Normal 분석", "확인 완료", "정상인데 위험처럼 보이는 8건 식별\nEvent 35/48: 구조적 어려운 정상", AMBER),
    ("Fold 안정성", "Pass", "compact13 Jaccard 34~40%\nFold 간 feature 선택 안정", GREEN),
    ("확장 학습 검증", "Pass", "weak positive 12건 + 후보 normal 70건\n확장 후 BA +2.1% 개선", GREEN),
]
for i, (title, status, desc, color) in enumerate(checks):
    col = i % 2
    row = i // 2
    cx = 0.5 + col * 4.7
    cy = 1.4 + row * 1.85
    card = add_card(slide, cx, cy, 4.4, 1.6, fill=WHITE)
    # Status badge
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + 3.2), Inches(cy + 0.15), Inches(1.0), Inches(0.25))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    set_rounded_rect(badge, 50000)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, cx + 0.2, cy + 0.15, 2.8, 0.3, title, font_size=15, color=SLATE, bold=True)
    add_text_box(slide, cx + 0.2, cy + 0.55, 4.0, 0.9, desc, font_size=11, color=SLATE_MUTED)

add_slide_number(slide, 9)


# ========== SLIDE 10: Fault Gate ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.5, 0.25, 9.0, 0.5, "Fault Gate 확립 (NB 21)", font_size=30, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.7, 9.0, 0.3,
             "compact13 + RandomForest (depth=3) + threshold 0.5", font_size=13, color=TEAL)

# Confusion matrix
cm_x, cm_y = 0.7, 1.3
# Header labels
add_text_box(slide, cm_x + 1.8, cm_y - 0.35, 3.5, 0.3, "모델 예측", font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.CENTER, bold=True)
add_text_box(slide, cm_x + 1.6, cm_y, 1.5, 0.4, "위험", font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.CENTER)
add_text_box(slide, cm_x + 3.1, cm_y, 1.5, 0.4, "정상", font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.CENTER)

# Row labels
add_text_box(slide, cm_x - 0.5, cm_y + 0.5, 1.6, 0.8, "실제\n장애", font_size=12,
             color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, bold=True)
add_text_box(slide, cm_x - 0.5, cm_y + 1.5, 1.6, 0.8, "실제\n정상", font_size=12,
             color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, bold=True)

# Cells
cm_cells = [
    (cm_x + 1.6, cm_y + 0.45, "49", "TP", GREEN),
    (cm_x + 3.1, cm_y + 0.45, "6", "FN", RGBColor(0x7F, 0x1D, 0x1D)),
    (cm_x + 1.6, cm_y + 1.45, "7", "FP", AMBER),
    (cm_x + 3.1, cm_y + 1.45, "28", "TN", DEEP_BLUE),
]
for (cx, cy, num, label, color) in cm_cells:
    cell = add_card(slide, cx, cy, 1.4, 0.9, fill=RGBColor(0x0E, 0x34, 0x60), shadow=False, rounded=True)
    add_text_box(slide, cx + 0.1, cy + 0.05, 1.2, 0.5, num, font_size=32, color=color, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.1, cy + 0.55, 1.2, 0.3, label, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8),
                 align=PP_ALIGN.CENTER)

# Right side metrics
card = add_card(slide, 5.8, 1.3, 3.7, 2.5, fill=RGBColor(0x0E, 0x34, 0x60), shadow=True, rounded=True)
metrics = [
    ("Recall", "0.891", "55건 중 49건 탐지"),
    ("FPR", "0.200", "35건 중 7건 오탐"),
    ("BA", "0.845", "균형 정확도"),
]
for i, (name, val, desc) in enumerate(metrics):
    my = 1.45 + i * 0.75
    add_text_box(slide, 6.0, my, 1.2, 0.3, name, font_size=12, color=RGBColor(0x94, 0xA3, 0xB8), bold=True)
    add_text_box(slide, 7.2, my, 1.0, 0.3, val, font_size=20, color=TEAL, bold=True)
    add_text_box(slide, 6.0, my + 0.3, 3.3, 0.25, desc, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))

# Bottom interpretation
add_text_box(slide, 0.5, 4.2, 9.0, 0.6,
             '"장애 10건이 발생하면 9건은 미리 경고. 다만 정상 5건 중 1건은 잘못 울림"',
             font_size=15, color=LIGHT_TEAL, align=PP_ALIGN.CENTER, bold=True)

add_slide_number(slide, 10, dark=True)


# ========== SLIDE 11: 4-class attempt ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 2")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5,
             "다음 질문: 모든 이벤트를 분류할 수 있는가? (NB 16)", font_size=26, color=SLATE, bold=True)

# 4 label boxes
labels_4 = [
    ("Normal", "비교 기준", "정상 운영 구간", DEEP_BLUE),
    ("Fault", "장애 위험", "사전 탐지 대상", RED),
    ("Task", "작업/정비", "정비 맥락 이벤트", AMBER),
    ("Activity", "활동 신호", "운영 활동 이벤트", TEAL),
]
for i, (name, role, desc, color) in enumerate(labels_4):
    cx = 0.5 + i * 2.35
    card = add_card(slide, cx, 1.3, 2.1, 1.8, fill=WHITE)
    # Color circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.75), Inches(1.45), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = name[0]
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, cx + 0.15, 2.15, 1.8, 0.25, name, font_size=14, color=SLATE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.15, 2.4, 1.8, 0.2, role, font_size=11, color=color, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.15, 2.65, 1.8, 0.3, desc, font_size=10, color=SLATE_MUTED,
                 align=PP_ALIGN.CENTER)

# Problem discovery box
card = add_card(slide, 0.5, 3.4, 9.0, 1.6, fill=RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, 0.7, 3.55, 8.6, 0.3, "문제 발견", font_size=16, color=RED, bold=True)
add_bullet_list(slide, 0.7, 3.9, 8.6, 1.0, [
    "4분류 단일 모델: Logistic BA 0.45 — 사실상 사용 불가",
    "핵심 원인: 네 라벨의 '시간적 의미'가 서로 다르다",
    "fault는 '사전 위험', task는 '사후 맥락', activity는 '결측 패턴 의존' — 같은 분류 모델로 묶을 수 없다",
], font_size=12, color=SLATE)

add_slide_number(slide, 11)


# ========== SLIDE 12: Labels are different ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.5, 0.25, 9.0, 0.5,
             "핵심 발견: 라벨은 같은 종류가 아니다", font_size=28, color=WHITE, bold=True)

# 3-column comparison table
table_data = [
    ["", "Fault", "Task", "Activity"],
    ["시간적 의미", "사전 위험 신호", "사후 맥락", "불확실"],
    ["예측 가능?", "Yes", "No (overlap 0.40)", "의심 (missingness)"],
    ["운영 역할", "Predictive Gate", "Context Classifier", "보조 신호"],
    ["근거", "NB 19-21 검증 통과", "NB 22 pre=실패, post=통과", "NB 23 missingness=1.00"],
    ["결정", "예측 모델로 사용", "사후 분류로 전환", "보조 신호로 보류"],
]
tbl = add_table_to_slide(slide, 0.5, 1.0, 9.0, table_data, header_color=RGBColor(0x0E, 0x34, 0x60),
                          font_size=12)
# Adjust table cell colors for dark theme
table = tbl.table
for row_idx in range(1, len(table_data)):
    for col_idx in range(len(table_data[0])):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x27, 0x44) if row_idx % 2 == 1 else RGBColor(0x0E, 0x34, 0x60)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)

# Bottom insight
add_text_box(slide, 0.5, 4.2, 9.0, 0.6,
             '"같은 데이터에서 나온 라벨이라도, 시간적 의미가 다르면 같은 모델로 묶을 수 없다"',
             font_size=15, color=LIGHT_TEAL, align=PP_ALIGN.CENTER, bold=True)

add_slide_number(slide, 12, dark=True)


# ========== SLIDE 13: Task is not predictive ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 2")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5,
             "Task: 사후 맥락이지, 사전 예측이 아니다 (NB 22)", font_size=24, color=SLATE, bold=True)

# Comparison table
task_table = [
    ["후보", "BA", "Recall", "FPR", "Overlap", "판정"],
    ["task_pre_1d", "0.80", "0.77", "0.17", "0.40 (40%)", "Overlap 과다"],
    ["task_no_overlap", "0.65", "0.68", "0.37", "0.00", "성능 부족"],
    ["task_post_1d", "통과", "통과", "통과", "-", "사후 분류로 채택"],
]
add_table_to_slide(slide, 0.5, 1.2, 9.0, task_table, font_size=12)

# Explanation
card = add_card(slide, 0.5, 3.0, 9.0, 2.0, fill=WHITE)
add_text_box(slide, 0.7, 3.15, 8.6, 0.3, "왜 Task는 예측이 아닌가?", font_size=16, color=SLATE, bold=True)
add_bullet_list(slide, 0.7, 3.5, 8.6, 1.3, [
    "task_pre_1d: 겉보기 성능은 좋지만 40%가 사건 기간과 겹침 (데이터 누출 위험)",
    "task_no_overlap: overlap 제거 시 BA 0.65로 하락 — 진짜 사전 신호가 약함",
    "task_post_1d: 사후 데이터로는 잘 분류됨 → 본질은 '사후 맥락 분류'",
], font_size=12, color=SLATE)

add_text_box(slide, 0.7, 4.6, 8.6, 0.3,
             '결정: "task를 예측했다"고 말하면 안 됨. "정비 후 맥락을 분류했다"가 정확',
             font_size=12, color=TEAL, bold=True)

add_slide_number(slide, 13)


# ========== SLIDE 14: Activity trap ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 2")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5,
             "Activity의 함정: 완벽해 보였지만... (NB 23)", font_size=24, color=SLATE, bold=True)

# Two comparison cards
# Left: looks perfect
card = add_card(slide, 0.5, 1.2, 4.2, 2.5, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 0.7, 1.35, 3.8, 0.3, "activity_pre_1d", font_size=16, color=GREEN, bold=True)
add_text_box(slide, 0.7, 1.75, 3.8, 0.5, "BA 1.00  |  Recall 1.00  |  FPR 0.00",
             font_size=18, color=GREEN, bold=True)
add_text_box(slide, 0.7, 2.35, 3.8, 0.5, "완벽한 성능!", font_size=24, color=GREEN, bold=True,
             align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.9, 3.8, 0.4, "...정말 그럴까?", font_size=14, color=SLATE_MUTED,
             align=PP_ALIGN.CENTER)

# Right: the truth
card = add_card(slide, 5.3, 1.2, 4.2, 2.5, fill=RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, 5.5, 1.35, 3.8, 0.3, "missingness_only 모델", font_size=16, color=RED, bold=True)
add_text_box(slide, 5.5, 1.75, 3.8, 0.5, "BA 1.00  |  역시 완벽",
             font_size=18, color=RED, bold=True)
add_text_box(slide, 5.5, 2.35, 3.8, 0.8,
             '"센서값이 아니라\n어디가 비어 있는지만 봐도\n완벽하게 맞춘다"',
             font_size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)

# Bottom: what happens when we remove missingness
card = add_card(slide, 0.5, 4.0, 9.0, 1.2, fill=WHITE)
add_text_box(slide, 0.7, 4.1, 8.6, 0.3, "결측 feature 제거 시", font_size=14, color=SLATE, bold=True)
add_text_box(slide, 0.7, 4.4, 4.0, 0.3, "BA 1.00 → 0.74로 하락", font_size=16, color=RED, bold=True)
add_text_box(slide, 5.0, 4.4, 4.6, 0.6,
             "결정: activity는 예측 성공이라 말하면 안 됨.\n보조 신호(overlap/missingness-sensitive candidate)로만 사용",
             font_size=12, color=TEAL, bold=True)

add_slide_number(slide, 14)


# ========== SLIDE 15: Hierarchical Routing ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.5, 0.2, 9.0, 0.5, "해법: 계층형 운영 판단 정책 (NB 24)", font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.65, 9.0, 0.3, "4분류 모델이 아니라 우선순위 기반 라우팅", font_size=13, color=TEAL)

# Decision tree flow
levels = [
    ("1", "Fault Gate Positive?", "predictive_fault_risk", "56건", GREEN),
    ("2", "Task Context Positive?", "maintenance_context_event", "42건", AMBER),
    ("3", "Activity Candidate?", "activity_context_signal", "47건", TEAL),
    ("4", "None", "normal_or_monitor", "34건", SLATE_MUTED),
]
for i, (num, question, label, count, color) in enumerate(levels):
    cy = 1.1 + i * 1.0
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(cy + 0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Question
    add_text_box(slide, 1.3, cy, 3.5, 0.4, question, font_size=15, color=WHITE, bold=True)

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(cy + 0.1), Inches(0.4), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

    # Result card
    card = add_card(slide, 5.4, cy - 0.05, 3.0, 0.55, fill=RGBColor(0x0E, 0x34, 0x60), shadow=False, rounded=True)
    add_text_box(slide, 5.55, cy, 2.1, 0.4, label, font_size=11, color=color, bold=True)
    add_text_box(slide, 7.6, cy, 0.7, 0.4, count, font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
                 align=PP_ALIGN.RIGHT)

    # Vertical line connecting
    if i < 3:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.88), Inches(cy + 0.5), Inches(0.03), Inches(0.55))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x66)
        line.line.fill.background()

# Key rule
add_text_box(slide, 0.5, 5.0, 9.0, 0.4,
             "핵심 규칙: fault 우선. 장애 위험이 잡히면 다른 신호가 있어도 fault로 분류",
             font_size=13, color=LIGHT_TEAL, bold=True, align=PP_ALIGN.CENTER)

add_slide_number(slide, 15, dark=True)


# ========== SLIDE 16: Model search rejected ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 3")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5,
             "더 나은 모델을 찾았지만, 채택하지 않았다 (NB 25)", font_size=24, color=SLATE, bold=True)

# Comparison table
comp_table = [
    ["", "기존 Gate (NB 21)", "새 후보 (NB 25)"],
    ["Recall (장애 탐지)", "0.891", "0.945"],
    ["FPR (정상 오탐률)", "0.200", "0.286"],
    ["놓친 장애 (FN)", "6건", "3건"],
    ["잘못 울린 정상 (FP)", "7건", "10건"],
    ["BA", "0.845", "0.830"],
]
add_table_to_slide(slide, 0.5, 1.2, 9.0, comp_table, font_size=13)

# Trade-off explanation
card = add_card(slide, 0.5, 3.6, 9.0, 1.5, fill=WHITE)
add_text_box(slide, 0.7, 3.7, 8.6, 0.3, "왜 채택하지 않았는가?", font_size=16, color=SLATE, bold=True)
add_bullet_list(slide, 0.7, 4.0, 8.6, 0.9, [
    "장애 3건을 더 잡는 대신, 정상 3건을 더 잘못 울림",
    "FPR 0.286 → M1 운영 기준(최대 0.25) 초과",
    '핵심: "최고 점수 모델 ≠ 최적 운영 모델"',
], font_size=12, color=SLATE, bold_items={2})

add_slide_number(slide, 16)


# ========== SLIDE 17: Final Recipe ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 3")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "최종 Fault Gate Recipe (NB 27)", font_size=28, color=SLATE, bold=True)

# Spec card
card = add_card(slide, 0.5, 1.2, 4.5, 3.8, fill=WHITE)
add_text_box(slide, 0.7, 1.35, 4.1, 0.3, "모델 사양", font_size=16, color=SLATE, bold=True)

specs = [
    ("Target", "fault pre-event risk (efd_possible)"),
    ("Window", "report/fault 이전 7일"),
    ("Features", "compact13_overlap (13개)"),
    ("Model", "Logistic (class_weight=balanced)"),
    ("Threshold", "0.6"),
    ("Training", "fixed eval 49 + expanded 82"),
]
for i, (key, val) in enumerate(specs):
    sy = 1.75 + i * 0.45
    add_text_box(slide, 0.8, sy, 1.4, 0.3, key, font_size=11, color=SLATE_MUTED, bold=True)
    add_text_box(slide, 2.2, sy, 2.6, 0.35, val, font_size=12, color=SLATE)

# Metrics card
card = add_card(slide, 5.3, 1.2, 4.2, 2.2, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 5.5, 1.35, 3.8, 0.3, "성능 지표", font_size=16, color=GREEN, bold=True)

recipe_metrics = [
    ("BA", "0.85"),
    ("Recall", "0.786 (14건 중 11건)"),
    ("FPR", "0.086 (35건 중 3건 오탐)"),
    ("FP", "3건"),
    ("FN", "3건"),
]
for i, (name, val) in enumerate(recipe_metrics):
    my = 1.75 + i * 0.32
    add_text_box(slide, 5.6, my, 1.0, 0.25, name, font_size=11, color=SLATE_MUTED, bold=True)
    add_text_box(slide, 6.6, my, 2.8, 0.25, val, font_size=12, color=SLATE, bold=True)

# Quality checks
card = add_card(slide, 5.3, 3.6, 4.2, 1.4, fill=WHITE)
add_text_box(slide, 5.5, 3.7, 3.8, 0.3, "Quality Checks: All Passed", font_size=13, color=GREEN, bold=True)
add_bullet_list(slide, 5.5, 4.0, 3.8, 0.9, [
    "Leakage: None",
    "Group overlap: Zero",
    "Hard normal 35/48: FP 0",
    "13번 recipe 재현 성공",
], font_size=10, color=SLATE_MUTED)

add_slide_number(slide, 17)


# ========== SLIDE 18: Service perspective ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_section_tag(slide, "Phase 3")
add_text_box(slide, 0.5, 0.5, 9.0, 0.5, "서비스 관점: Agent 입력 구조", font_size=28, color=SLATE, bold=True)

# Schema table
schema_data = [
    ["필드", "타입", "의미"],
    ["fault_probability", "float", "장애 위험 점수 (0~1)"],
    ["fault_pred", "bool", "장애 gate 통과 여부"],
    ["task_context_probability", "float", "작업/정비 맥락 점수"],
    ["activity_candidate_prob", "float", "activity 보조 신호 점수"],
    ["final_operational_label", "string", "최종 운영 라벨"],
    ["supporting_context_tags", "list", "보조 맥락 태그"],
    ["review_required", "bool", "사람 확인 필요 여부"],
]
add_table_to_slide(slide, 0.5, 1.2, 6.5, schema_data, col_widths=[2.2, 0.8, 3.5], font_size=11)

# Side explanation
card = add_card(slide, 7.3, 1.2, 2.2, 3.0, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 7.45, 1.35, 1.9, 0.3, "핵심 포인트", font_size=13, color=GREEN, bold=True)
add_bullet_list(slide, 7.45, 1.7, 1.9, 2.2, [
    "정답 하나가 아닌 구조화된 입력",
    "운영자 판단 가능",
    "review_required는 약점이 아닌 장점",
    "불확실한 케이스를 숨기지 않음",
], font_size=10, color=SLATE)

add_slide_number(slide, 18)


# ========== SLIDE 19: Pros & Cons ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 9.0, 0.5, "장점과 한계", font_size=30, color=SLATE, bold=True)

# Strengths
card = add_card(slide, 0.5, 1.0, 4.2, 3.8, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text_box(slide, 0.7, 1.15, 3.8, 0.3, "장점", font_size=18, color=GREEN, bold=True)
strengths = [
    "해석 가능: 왜 fault인지, 왜 context인지 분리 설명 가능",
    "운영 기준 존재: recall만 보지 않고 FPR 제약을 둠",
    "라벨 의미 존중: task/activity를 억지로 예측 대상으로 만들지 않음",
]
add_bullet_list(slide, 0.7, 1.55, 3.8, 2.8, strengths, font_size=12, color=SLATE, space_after=10)

# Limitations
card = add_card(slide, 5.3, 1.0, 4.2, 3.8, fill=WHITE)
add_text_box(slide, 5.5, 1.15, 3.8, 0.3, "한계", font_size=18, color=SLATE_MUTED, bold=True)
limits = [
    "샘플 수 작음 (fault 기준 90건)",
    "activity 결측 패턴 의존성",
    "threshold 안정성 추가 검토 필요",
    "FPR 20%도 실제 운영에선 부담 가능",
    "기계실/기간 확장 시 결과 변동 가능",
]
add_bullet_list(slide, 5.5, 1.55, 3.8, 3.0, limits, font_size=12, color=SLATE, space_after=8)

add_slide_number(slide, 19)


# ========== SLIDE 20: Next Steps ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 9.0, 0.5, "다음 단계", font_size=30, color=SLATE, bold=True)

steps = [
    ("1", "Fault Gate Threshold 운영 검토", "0.45 / 0.50 / 0.55 / 0.60에서 알람 수와 미탐 수 비교", TEAL),
    ("2", "Activity Relabeling / Missingness Audit", "activity가 물리적 활동인지 데이터 수집 상태인지 분리", AMBER),
    ("3", "Task Context Feature 정리", "예측 모델 성능 개선 대상 아님 → 운영 설명 레이어로 유지", SLATE_MUTED),
    ("4", "Agent 입력 Schema 설계", "fault_probability, context_tags, review_required → Agent 입력 형태 정리", DEEP_BLUE),
    ("5", "운영자 피드백 루프", "FP 7건, FN 6건을 사람이 보면 threshold/feature 해석 조정 가능", GREEN),
]
for i, (num, title, desc, color) in enumerate(steps):
    cy = 1.0 + i * 0.85
    # Number
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(cy + 0.05), Inches(0.35), Inches(0.35))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    card = add_card(slide, 1.1, cy, 8.4, 0.7, fill=WHITE)
    add_text_box(slide, 1.3, cy + 0.05, 3.5, 0.3, title, font_size=14, color=SLATE, bold=True)
    add_text_box(slide, 1.3, cy + 0.35, 8.0, 0.3, desc, font_size=11, color=SLATE_MUTED)

add_slide_number(slide, 20)


# ========== SLIDE 21: Key Lessons ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.5, 0.25, 9.0, 0.5, "27개 모델에서 배운 것", font_size=30, color=WHITE, bold=True)

lessons = [
    ('"Feature를 늘리는 것보다\n핵심을 고르는 것이 낫다"',
     "70→154: +2.8%  |  154→13: +18%", TEAL),
    ('"모든 라벨이 같은 종류의\n예측 대상은 아니다"',
     "fault=사전위험  |  task=사후맥락  |  activity=결측의존", AMBER),
    ('"최고 점수 모델이\n최적 운영 모델은 아니다"',
     "Recall 0.945 + FPR 0.286 → 기각  |  운영 기준이 우선", GREEN),
]
for i, (quote, evidence, color) in enumerate(lessons):
    cy = 1.0 + i * 1.45
    card = add_card(slide, 0.8, cy, 8.4, 1.2, fill=RGBColor(0x0E, 0x34, 0x60), shadow=True, rounded=True)
    # Number
    add_text_box(slide, 1.0, cy + 0.15, 0.5, 0.4, str(i + 1), font_size=28, color=color, bold=True)
    add_text_box(slide, 1.5, cy + 0.1, 6.5, 0.7, quote, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, 1.5, cy + 0.8, 7.0, 0.3, evidence, font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

add_slide_number(slide, 21, dark=True)


# ========== SLIDE 22: Closing ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide_bg(slide)

add_text_box(slide, 0.5, 1.5, 9.0, 1.0,
             "M1은 장애 위험을 먼저 잡고,\n맥락을 덧붙이는\n계층형 운영 판단 정책이다.",
             font_size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Teal accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.2), Inches(3.0), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = TEAL
shape.line.fill.background()

add_text_box(slide, 0.5, 3.5, 9.0, 0.6,
             "fault gate: 기존 모델 유지  |  task: 사후 맥락  |  activity: 보조 신호",
             font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 4.4, 9.0, 0.5, "Q & A", font_size=24, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)

add_slide_number(slide, 22, dark=True)


# ========== SAVE ==========
output_path = os.path.join(r"C:\Users\flfkr\OneDrive\바탕 화면\HeatGrid_re",
                            "M1_모델개발여정_발표자료.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
